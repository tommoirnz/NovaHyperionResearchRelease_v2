"""
NOVA HYPERION — Unified Web Interface
Merged Features:
  * Web TTS (edge_tts) + mute button + AMBIENT HUM (Web Audio)
  * Voice input (Whisper)
  * MathJax rendering + marked.js markdown
  * Code-block copy buttons
  * Inline image / diagram / audio / video
  * Save-response button + download/new tab
  * Drag-and-drop file upload (full preview)
  * Online/offline status dot
  * PARALLAX STARFIELD (mouse-tracked, shooting stars)
  * HEALTH CONDUIT animation
  * WARP CORE typing indicator
  * LCARS animated borders
"""

import http.server
import socketserver
import json
import threading
import webbrowser
import os
from urllib.parse import urlparse
from datetime import datetime
from personality_manager import personality_manager
_history_lock = threading.Lock()


# =============================================================================
#  HTTP HANDLER
# =============================================================================

class NovaWebHandler(http.server.SimpleHTTPRequestHandler):
    nova = None

    def do_GET(self):
        path = urlparse(self.path).path

        if path == '/':
            self._serve_html()
        elif path == '/api/history':
            self._serve_json(self.nova.conversation_history)
        elif path == '/api/state':
            self._serve_json({
                'model': self.nova.ai.model,
                'thinking': getattr(self.nova, '_thinking', False),
                'tts': getattr(self.nova, '_tts_on', False),
                'web_tts': getattr(self.nova, '_web_tts_on', True),
                'ambient_hum': getattr(self.nova, '_ambient_hum_on', True),
            })
        elif path == '/api/ping':
            self._serve_json({'status': 'ok', 'timestamp': datetime.now().isoformat()})
        elif path == '/api/personalities':
            active = personality_manager.active_name
            data = [{"name": p["name"], "description": p.get("description", ""),
                     "temperature": p.get("temperature", 0.7), "voice": p.get("voice", {})}
                    for p in personality_manager.all]
            self._serve_json({"personalities": data, "active": active})

        elif path.startswith('/api/stream'):
            from urllib.parse import parse_qs
            params = parse_qs(urlparse(self.path).query)
            filepath = params.get('file', [''])[0]
            if filepath and os.path.exists(filepath):
                ext = os.path.splitext(filepath)[1].lower()
                mime = {'.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.ogg': 'audio/ogg',
                        '.mp4': 'video/mp4', '.webm': 'video/webm', '.m4a': 'audio/mp4',
                        '.flac': 'audio/flac'}.get(ext, 'application/octet-stream')
                size = os.path.getsize(filepath)
                rng = self.headers.get('Range')
                try:
                    if rng:
                        s, e = rng.replace('bytes=', '').split('-')
                        s = int(s)
                        e = int(e) if e else size - 1
                        self.send_response(206)
                        self.send_header('Content-Type', mime)
                        self.send_header('Content-Range', f'bytes {s}-{e}/{size}')
                        self.send_header('Content-Length', e - s + 1)
                        self.send_header('Accept-Ranges', 'bytes')
                        self.end_headers()
                        with open(filepath, 'rb') as f:
                            f.seek(s)
                            self.wfile.write(f.read(e - s + 1))
                    else:
                        self.send_response(200)
                        self.send_header('Content-Type', mime)
                        self.send_header('Content-Length', size)
                        self.send_header('Accept-Ranges', 'bytes')
                        self.end_headers()
                        with open(filepath, 'rb') as f:
                            self.wfile.write(f.read())
                except Exception:
                    pass
            else:
                self.send_error(404)

        elif path.startswith('/images/'):
            filename = path[8:]
            filepath = os.path.join(os.getcwd(), 'web_images', filename)
            if os.path.exists(filepath):
                self.send_response(200)
                ct = ('image/png' if filename.endswith('.png') else
                      'image/jpeg' if filename.lower().endswith(('.jpg', '.jpeg')) else
                      'image/svg+xml' if filename.endswith('.svg') else
                      'application/octet-stream')
                self.send_header('Content-type', ct)
                self.end_headers()
                with open(filepath, 'rb') as f:
                    self.wfile.write(f.read())
            else:
                self.send_error(404)

        elif path.startswith('/plots/'):
            filename = path[7:]
            filepath = os.path.join(os.getcwd(), 'plots', filename)
            if os.path.exists(filepath):
                self.send_response(200)
                self.send_header('Content-type', 'text/html; charset=utf-8')
                self.send_header('Cache-Control', 'no-cache')
                self.end_headers()
                with open(filepath, 'rb') as f:
                    self.wfile.write(f.read())
            else:
                self.send_error(404)

        else:
            self.send_error(404)

    def do_OPTIONS(self):
        self.send_response(204)
        self._add_cors_headers()
        self.end_headers()

    def do_POST(self):
        path = urlparse(self.path).path
        if path == '/api/send':
            self._handle_send()
        elif path == '/api/voice':
            self._handle_voice()
        elif path == '/api/speak':
            self._handle_speak()
        elif path == '/api/imagine':
            self._handle_imagine()
        elif path == '/api/clear':
            self.nova.root.after(0, lambda: self.nova._new_chat())
            self._serve_json({'status': 'cleared'})
        elif path == '/api/upload':
            self._handle_upload()
        elif path == '/api/tts':
            self._handle_tts_toggle()
        elif path == '/api/ambient':
            self._handle_ambient_toggle()
        elif path == '/api/stop_speaking':
            self.nova.root.after(0, self.nova._stop_speaking)
            self._serve_json({'status': 'stopped'})


        elif path == '/api/personalities/set':
            length = int(self.headers['Content-Length'])
            body = json.loads(self.rfile.read(length).decode('utf-8'))
            name = body.get('name')
            if name is None:
                personality_manager.deactivate()
                # Restore default temperature
                if self.nova and hasattr(self.nova, 'ai'):
                    self.nova.ai.current_temperature = 0.3
                self._serve_json({'ok': True, 'active': None})
            else:
                try:
                    p = personality_manager.activate(name)
                    # Actually apply the temperature to the AI engine
                    if self.nova and hasattr(self.nova, 'ai'):
                        self.nova.ai.current_temperature = p.get('temperature', 0.3)
                        self.nova.log(f"[PERSONA] {name} → temp={p.get('temperature', 0.3)}")
                    self._serve_json({'ok': True, 'active': name,
                                      'temperature': p['temperature'], 'voice': p['voice']})
                except KeyError as e:
                    self._serve_json({'ok': False, 'error': str(e)})

        else:
            self.send_error(404)

    def _add_cors_headers(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type, Authorization')
        self.send_header('Access-Control-Max-Age', '86400')

    def _serve_html(self):
        self.send_response(200)
        self.send_header('Content-type', 'text/html; charset=utf-8')
        self.send_header('Cache-Control', 'no-cache, no-store, must-revalidate')
        self.send_header('Content-Security-Policy',
                         "default-src 'self'; "
                         "script-src 'self' 'unsafe-inline' https://cdn.jsdelivr.net https://cdnjs.cloudflare.com; "
                         "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
                         "font-src 'self' https://fonts.gstatic.com https://cdn.jsdelivr.net; "
                         "img-src 'self' data: blob:; media-src 'self' blob: data:; "
                         "frame-src 'self'; "
                         "connect-src 'self'; frame-ancestors 'none';")
        self.end_headers()
        self.wfile.write(HTML_TEMPLATE.encode('utf-8'))

    def _serve_json(self, data):
        self.send_response(200)
        self.send_header('Content-type', 'application/json')
        self._add_cors_headers()
        self.end_headers()

        def _s(o):
            if hasattr(o, 'strftime'): return o.strftime('%H:%M')
            return str(o)

        self.wfile.write(json.dumps(data, default=_s).encode('utf-8'))

    def _handle_voice(self):
        try:
            length = int(self.headers['Content-Length'])
            audio_data = self.rfile.read(length)
            self.nova.log(f"[VOICE] Received {len(audio_data)} bytes")
            if not audio_data:
                self._serve_json({'status': 'ok', 'transcript': ''});
                return
            whisper = getattr(self.nova, 'whisper', None)
            if not whisper or not whisper.model_loaded:
                self._serve_json({'error': 'Whisper not ready'});
                return
            try:
                from pydub import AudioSegment
                import soundfile as sf, io, numpy as np
                seg = AudioSegment.from_file(io.BytesIO(audio_data))
                seg = seg.set_frame_rate(16000).set_channels(1)
                self.nova.log(f"[VOICE] {seg.duration_seconds:.1f}s")
                wav_io = io.BytesIO();
                seg.export(wav_io, format='wav');
                wav_io.seek(0)
                samples, _ = sf.read(wav_io, dtype='float32')
                text = whisper.asr_model.transcribe(samples, 16000)
                self.nova.log(f"[VOICE] '{text}'")
                self._serve_json({'status': 'ok', 'transcript': text.strip() if text else ''})
            except ImportError:
                self._serve_json({'error': 'pydub not installed'})
            except Exception as e:
                self._serve_json({'error': str(e)})
        except Exception as e:
            self._serve_json({'error': str(e)})

    def _handle_speak(self):
        try:
            length = int(self.headers['Content-Length'])
            data = json.loads(self.rfile.read(length).decode('utf-8'))
            text = data.get('text', '').strip()
            if not text or not getattr(self.nova, '_web_tts_on', True) \
                    or getattr(self.nova, '_tts_on', False):
                self.send_response(204);
                self.end_headers();
                return

            import asyncio, edge_tts, re

            text = re.sub(
                u'[\U0001F300-\U0001FFFF'
                u'\u2600-\u26FF'
                u'\u2700-\u27BF'
                u'\u2300-\u23FF'
                u'\u25A0-\u25FF'
                u'\u2B00-\u2BFF'
                u']+',
                ' ', text, flags=re.UNICODE
            )

            math_speech = getattr(self.nova, 'math_speech', None)
            if math_speech:
                clean = math_speech.make_speakable_text(text, speak_math=True)
            else:
                clean = re.sub(r'\$\$[\s\S]+?\$\$', ' ', text)
                clean = re.sub(r'\$[^\$\n]{1,300}?\$', ' ', clean)
                clean = re.sub(r'\\\[[\s\S]+?\\\]', ' ', clean)
                clean = re.sub(r'\\\(.+?\\\)', ' ', clean)
                clean = re.sub(r'```[\s\S]*?```', ' ', clean)
                clean = re.sub(r'`[^`]+`', ' ', clean)
                clean = re.sub(r'^#{1,6}\s+(.+)$', r'\1.', clean, flags=re.MULTILINE)
                clean = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', clean)
                clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)
                clean = re.sub(r'^[-*]\s+', '', clean, flags=re.MULTILINE)
                clean = re.sub(r'---+', '', clean)
                clean = re.sub(r'\n{2,}', '. ', clean)
                clean = re.sub(r'[^\w\s\.,!?;:\-\(\)\/\+\=%"\']+', ' ', clean)

            clean = re.sub(r'\s+', ' ', clean).strip()
            if not clean:
                self.send_response(204);
                self.end_headers();
                return


            # Personality voice takes priority
            p_voice = personality_manager.active_voice()
            voice = p_voice.get('edge_voice', 'en-AU-NatashaNeural')
            # Fall back to UI combo if no personality active
            if not personality_manager.active:
                engine_obj = getattr(self.nova, 'tts_engine_combo', None)
                engine_str = engine_obj.get() if engine_obj else 'sapi5'
                if engine_str == 'edge':
                    vc = getattr(self.nova, 'edge_voice_combo', None)
                    voice = vc.get() if vc else 'en-AU-NatashaNeural'
                else:
                    voice = 'en-GB-SoniaNeural'

            # Build rate string
            raw_rate = p_voice.get('speech_rate', 0)
            if isinstance(raw_rate, str):
                rate_str = raw_rate  # already "-20%" etc
            else:
                rate_str = f"{raw_rate * 10:+d}%"  # integer → "+0%"

            print(f"[WEB TTS DEBUG] voice='{voice}' rate_str='{rate_str}'")

            _clean = clean  # capture for closure

            async def _run():
                comm = edge_tts.Communicate(_clean, voice, rate=rate_str)  # type: ignore
                chunks = []
                async for chunk in comm.stream():  # type: ignore
                    if chunk['type'] == 'audio':
                        chunks.append(chunk['data'])
                return b''.join(chunks)

            loop = asyncio.new_event_loop()
            audio = loop.run_until_complete(_run())
            loop.close()

            if getattr(self.nova, '_tts_recording', False):
                from datetime import datetime
                import os
                try:
                    rec_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "recordings")
                    os.makedirs(rec_dir, exist_ok=True)
                    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                    path = os.path.join(rec_dir, f"nova_web_{ts}.mp3")
                    with open(path, "wb") as f:
                        f.write(audio)
                    self.nova.log(f"[WEB TTS] Saved → {path}")
                except Exception as rec_e:
                    self.nova.log(f"[WEB TTS] Recording failed: {rec_e}")

            self.send_response(200)
            self.send_header('Content-Type', 'audio/mpeg')
            self.send_header('Content-Length', len(audio))
            self.end_headers()
            self.wfile.write(audio)

        except Exception as e:
            self.nova.log(f"[SPEAK] {e}")
            self.send_response(204);
            self.end_headers()

    def _handle_ambient_toggle(self):
        current = getattr(self.nova, '_ambient_hum_on', True)
        self.nova._ambient_hum_on = not current
        self.nova.log(f"[AMBIENT] {'Enabled' if self.nova._ambient_hum_on else 'Disabled'}")
        self._serve_json({'ambient': self.nova._ambient_hum_on})

    def _handle_send(self):
        length = int(self.headers['Content-Length'])
        data = json.loads(self.rfile.read(length).decode('utf-8'))
        msg = data.get('message', '')
        if msg:
            self.nova.root.after(0, lambda: self._send_to_nova(msg))
            self._serve_json({'status': 'processing'})

    def _send_to_nova(self, msg):
        if '[IMAGE_FILE:' in msg:
            import re
            match = re.search(r'\[IMAGE_FILE:\s*(.+?)\]', msg)
            if match:
                img_path = match.group(1).strip()
                if os.path.exists(img_path):
                    self.nova.log(f"[UPLOAD] Image file detected → vision: {img_path}")
                    comment = msg.split('User comment:')[-1].strip() if 'User comment:' in msg else ''
                    clean_msg = comment or 'Describe everything you can see in this image in detail.'

                    # ── PATCH: nudge Nova to offer to solve equations/problems ──
                    if not comment:
                        clean_msg = (
                            'Describe everything you can see in this image in detail. '
                            'If the image contains a mathematical equation, formula, physics problem, '
                            'circuit diagram, or any technical problem, describe it clearly and offer '
                            'to solve or analyse it.'
                        )
                    self.nova._append_conv("user", clean_msg)
                    self.nova.conversation_history.append({"role": "user", "content": clean_msg})
                    self.nova._thinking = True

                    def _vision_task():
                        try:
                            self.nova.log(f"[VISION] model={self.nova.ai.model}, img={img_path}")
                            result = self.nova.ai.generate(clean_msg, image_path=img_path, use_planning=False)
                            reply = result or "No response from model."
                            self.nova._append_conv("assistant", reply)
                            self.nova.conversation_history.append({"role": "assistant", "content": reply})
                            self.nova.state["last_task"] = clean_msg  # ← Fix 1
                            self.nova._save_history(clean_msg, reply)
                            self.nova._deliver_tool_result(reply)
                        except Exception as e:
                            self.nova.log(f"[VISION ERROR] {e}")
                            self.nova._deliver_tool_result(f"Vision error: {e}")
                        finally:
                            self.nova._thinking = False

        if "I've attached" in msg and "```" in msg:
            try:
                import re
                blocks = re.findall(
                    r'\*\*File: (.+?)\*\*.*?Content:\s*```\r?\n(.*?)\n```', msg, re.DOTALL)
                if blocks:
                    comment = msg.split('User comment:')[-1].strip() if 'User comment:' in msg else ''
                    parts = [f'Content of "{fn}":\n```\n{fc.strip()}\n```' for fn, fc in blocks]
                    clean = '\n\n'.join(parts) + '\n\n' + (comment or 'Please summarise this content.')
                    self.nova._append_conv("user", clean)
                    self.nova.conversation_history.append({"role": "user", "content": clean})
                    self.nova._thinking = True
                    threading.Thread(target=self.nova._process_input, args=(clean,), daemon=True).start()
                    return
            except Exception as e:
                self.nova.log(f"[UPLOAD] extraction failed: {e}")
        self.nova._append_conv("user", msg)
        self.nova.conversation_history.append({"role": "user", "content": msg})
        self.nova._thinking = True
        threading.Thread(target=self.nova._process_input, args=(msg,), daemon=True).start()

    def _handle_imagine(self):
        length = int(self.headers['Content-Length'])
        data = json.loads(self.rfile.read(length).decode('utf-8'))
        msg = data.get('message', '')
        if msg:
            prompt = (f"You are in an imaginative, free-thinking mode.\n"
                      f"Approach this with maximum creativity and lateral thinking.\n"
                      f"User prompt: {msg}")
            self.nova.root.after(0, lambda: self._imagine_to_nova(prompt, msg))
            self._serve_json({'status': 'processing'})

    def _imagine_to_nova(self, prompt, original):
        # Check for image in the original payload before wrapping
        if '[IMAGE_FILE:' in original:
            import re
            match = re.search(r'\[IMAGE_FILE:\s*(.+?)\]', original)
            if match:
                img_path = match.group(1).strip()
                if os.path.exists(img_path):
                    comment = original.split('User comment:')[-1].strip() \
                        if 'User comment:' in original else \
                        'Describe this image with maximum creativity and lateral thinking.'
                    self.nova._append_conv("user", f"✨ [Holodeck] {comment}")
                    self.nova.conversation_history.append({"role": "user", "content": f"✨ [Holodeck] {comment}"})
                    self.nova._thinking = True

                    def _vision_task():
                        try:
                            imagine_prompt = (f"You are in an imaginative, free-thinking mode.\n"
                                              f"Approach this with maximum creativity.\n"
                                              f"User prompt: {comment}")
                            result = self.nova.ai.generate(imagine_prompt, image_path=img_path, use_planning=False)
                            reply = result or "No response."
                            self.nova._append_conv("assistant", reply)
                            self.nova.conversation_history.append({"role": "assistant", "content": reply})
                            self.nova.state["last_task"] = comment  # ← Fix 1
                            self.nova._save_history(comment, reply)
                            self.nova._deliver_tool_result(reply)
                        except Exception as e:
                            self.nova.log(f"[HOLODECK VISION ERROR] {e}")
                            self.nova._deliver_tool_result(f"Vision error: {e}")
                        finally:
                            self.nova._thinking = False

        # Normal holodeck path (no image)
        self.nova._append_conv("user", f"✨ [Holodeck] {original}")
        self.nova.conversation_history.append({"role": "user", "content": f"✨ [Holodeck] {original}"})
        self.nova._thinking = True

        def _holodeck_task():
            old_temp = self.nova.ai.current_temperature
            try:
                self.nova.ai.current_temperature = 0.9
                self.nova.log(f"[HOLODECK] Temperature → 0.9 (was {old_temp})")
                self.nova._process_input(original)  # ← only change from your old version
            except Exception as e:
                self.nova.log(f"[HOLODECK ERROR] {e}")
                self.nova._deliver_tool_result(f"Holodeck error: {e}")
            finally:
                self.nova.ai.current_temperature = old_temp
                self.nova.log(f"[HOLODECK] Temperature restored → {old_temp}")
                self.nova._thinking = False

        threading.Thread(target=_holodeck_task, daemon=True).start()
    def _handle_tts_toggle(self):
        current = getattr(self.nova, '_web_tts_on', True)
        self.nova._web_tts_on = not current
        if not self.nova._web_tts_on:
            self.nova.root.after(0, self.nova._stop_speaking)
        else:
            self.nova._tts_stop = False
        self.nova.log(f"[WEB TTS] {'Enabled' if self.nova._web_tts_on else 'Disabled'}")
        self._serve_json({'tts': self.nova._web_tts_on})

    def _handle_upload(self):
        try:
            MAX_FILE = 100 * 1024 * 1024
            MAX_TOTAL = 500 * 1024 * 1024
            MAX_PREVIEW = 50_000
            MAX_CODE = 100_000

            ct = self.headers.get('Content-Type', '')
            cl = int(self.headers.get('Content-Length', 0))
            if cl > MAX_TOTAL:
                self.send_error(413, 'Upload too large');
                return
            if 'multipart/form-data' not in ct:
                self.send_error(400, 'Expected multipart/form-data');
                return

            boundary = None
            for part in ct.split(';'):
                p = part.strip()
                if p.startswith('boundary='):
                    boundary = p[9:].strip().encode('utf-8');
                    break
            if not boundary:
                self.send_error(400, 'Missing boundary');
                return

            raw = self.rfile.read(cl)
            files = {}
            for part in raw.split(b'--' + boundary)[1:]:
                if part in (b'--\r\n', b'--', b''): continue
                sep = b'\r\n\r\n' if b'\r\n\r\n' in part else b'\n\n'
                if sep not in part: continue
                hdr, body = part.split(sep, 1)
                if body.endswith(b'\r\n'): body = body[:-2]
                disp = {}
                for line in hdr.decode('utf-8', 'replace').splitlines():
                    if 'Content-Disposition' in line:
                        for tok in line.split(';'):
                            tok = tok.strip()
                            if '=' in tok:
                                k, v = tok.split('=', 1)
                                disp[k.strip()] = v.strip().strip('"')
                if disp.get('filename'):
                    files[disp.get('name', '')] = {'filename': disp['filename'], 'data': body}

            if 'files' not in files:
                self.send_error(400, 'No file field');
                return

            obj = files['files']
            items = obj if isinstance(obj, list) else [obj]
            for fi in items:
                if len(fi['data']) > MAX_FILE:
                    self.send_error(413, f'{fi["filename"]} too large');
                    return

            upload_dir = os.path.join(os.getcwd(), 'uploads')
            os.makedirs(upload_dir, exist_ok=True)
            results = []

            IMAGE = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
            CODE = {'.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.rs', '.go', '.rb', '.php'}
            TEXT = {'.txt', '.html', '.css', '.json', '.md', '.csv', '.xml', '.yaml', '.yml',
                    '.ini', '.toml', '.bat', '.sh', '.log', '.conf', '.cfg'}

            for fi in items:
                filename = os.path.basename(fi['filename'])
                file_data = fi['data']
                name, ext = os.path.splitext(filename)
                extl = ext.lower()

                save_path = os.path.join(upload_dir, filename)
                n = 1
                while os.path.exists(save_path):
                    save_path = os.path.join(upload_dir, f"{name}_{n}{ext}");
                    n += 1
                with open(save_path, 'wb') as f:
                    f.write(file_data)

                if extl in IMAGE:
                    preview = f"[IMAGE_FILE: {save_path}]"
                elif extl in CODE:
                    try:
                        t = file_data.decode('utf-8', 'replace')
                        preview = t[:MAX_CODE] + ('\n...[truncated]' if len(t) > MAX_CODE else '')
                    except:
                        preview = f"[Cannot decode: {filename}]"
                elif extl in TEXT:
                    try:
                        t = file_data.decode('utf-8', 'replace')
                        preview = t[:MAX_PREVIEW] + ('\n...[truncated]' if len(t) > MAX_PREVIEW else '')
                    except:
                        preview = f"[Cannot decode: {filename}]"
                elif extl in ('.xlsx', '.xls'):
                    try:
                        import pandas as pd, io as _io
                        xl = pd.ExcelFile(_io.BytesIO(file_data))
                        sheets = xl.sheet_names
                        df = pd.read_excel(_io.BytesIO(file_data), sheet_name=sheets[0])
                        preview = f"Excel: {len(sheets)} sheet(s): {', '.join(sheets[:5])}\n"
                        preview += f"Shape: {df.shape[0]}×{df.shape[1]}\n{df.head(5).to_string()}"
                        if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'
                    except ImportError:
                        preview = "[Excel: pandas not installed]"
                    except Exception as e:
                        preview = f"[Excel error: {e}]"
                elif extl == '.pdf':
                    try:
                        from pypdf import PdfReader
                        import io as _io
                        reader = PdfReader(_io.BytesIO(file_data))
                        preview = f"PDF: {len(reader.pages)} pages\n\n"
                        total = 0
                        for i, pg in enumerate(reader.pages[:3]):
                            t = ' '.join((pg.extract_text() or '').split())
                            if total + len(t) > MAX_PREVIEW: break
                            preview += f"--- Page {i + 1} ---\n{t}\n\n";
                            total += len(t)
                    except ImportError:
                        preview = "[PDF: pypdf not installed]"
                    except Exception as e:
                        preview = f"[PDF error: {e}]"
                elif extl in ('.pptx', '.ppt'):
                    try:
                        from pptx import Presentation
                        import io as _io
                        prs = Presentation(_io.BytesIO(file_data))
                        preview = f"PowerPoint: {len(prs.slides)} slides\n\n"
                        for i, slide in enumerate(prs.slides[:5]):
                            txts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text.strip()]
                            preview += f"--- Slide {i + 1} ---\n" + '\n'.join(txts) + '\n\n'
                            if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'; break
                    except ImportError:
                        preview = "[PPTX: python-pptx not installed]"
                    except Exception as e:
                        preview = f"[PPTX error: {e}]"
                elif extl in ('.docx', '.doc'):
                    try:
                        from docx import Document
                        import io as _io
                        doc = Document(_io.BytesIO(file_data))
                        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                        preview = '\n\n'.join(paras[:50])
                        if len(preview) > MAX_PREVIEW: preview = preview[:MAX_PREVIEW] + '\n...[truncated]'
                    except ImportError:
                        preview = "[DOCX: python-docx not installed]"
                    except Exception as e:
                        preview = f"[DOCX error: {e}]"
                else:
                    preview = f"[Binary: {filename}, {len(file_data) / (1024 * 1024):.2f} MB]"

                results.append({
                    'filename': os.path.basename(save_path),
                    'original_name': filename,
                    'size': len(file_data),
                    'path': save_path,
                    'preview': preview,
                })
                self.nova.log(f"[UPLOAD] {filename} ({len(file_data)} bytes)")

            self._serve_json({'status': 'success', 'files': results})
        except Exception as e:
            self.nova.log(f"[UPLOAD ERROR] {e}")
            self.send_error(500, str(e))

    def log_message(self, fmt, *args):
        msg = fmt % args
        if any(p in msg for p in ['/api/history', '/api/state', '/api/ping']): return
        if '"GET /api/' in msg and ('200' in msg or '304' in msg): return
        if ('/images/' in msg or '404' in msg or '500' in msg) and self.nova:
            self.nova.log(f"[WEB] {msg}")


# =============================================================================
#  SERVER WRAPPER
# =============================================================================

class NovaWebServer:
    def __init__(self, nova_app, port=8080, bind_all=False):
        self.nova = nova_app
        self.port = port
        self.bind_all = bind_all
        self.server = None
        self.thread = None

    def start(self):
        if self.server: return
        NovaWebHandler.nova = self.nova
        script_dir = os.path.dirname(os.path.abspath(__file__))
        cert = os.path.join(script_dir, 'cert.pem')
        key = os.path.join(script_dir, 'key.pem')
        use_ssl = os.path.exists(cert) and os.path.exists(key)
        host = "0.0.0.0" if self.bind_all else "127.0.0.1"

        if use_ssl:
            import ssl
            ctx = ssl.SSLContext(ssl.PROTOCOL_TLS_SERVER)
            ctx.load_cert_chain(cert, key)

            class SSLTCPServer(socketserver.ThreadingTCPServer):
                def get_request(self_):
                    ns, fa = self_.socket.accept()
                    return ctx.wrap_socket(ns, server_side=True), fa

            self.server = SSLTCPServer((host, self.port), NovaWebHandler)
            self.nova.log("[WEB] SSL enabled (HTTPS)");
            protocol = "https"
        else:
            self.server = socketserver.ThreadingTCPServer((host, self.port), NovaWebHandler)
            self.nova.log("[WEB] HTTP mode");
            protocol = "http"

        self.server.allow_reuse_address = True
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()
        ip = self._get_local_ip() if self.bind_all else "127.0.0.1"
        self.nova.log(f"[WEB] {protocol}://{ip}:{self.port}")

    def _get_local_ip(self):
        import socket
        try:
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.connect(("8.8.8.8", 80));
            ip = s.getsockname()[0];
            s.close();
            return ip
        except:
            return "localhost"

    def stop(self):
        if self.server:
            self.server.shutdown();
            self.server.server_close()
            self.server = None;
            self.nova.log("[WEB] Server stopped")

    def open_browser(self):
        proto = "https" if self.bind_all else "http"
        ip = self._get_local_ip() if self.bind_all else "127.0.0.1"
        webbrowser.open(f"{proto}://{ip}:{self.port}")


# =============================================================================
#  HTML TEMPLATE (Merged with Tkinter visual elements)
# =============================================================================

HTML_TEMPLATE = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, viewport-fit=cover, interactive-widget=resizes-content">
<meta name="apple-mobile-web-app-capable" content="yes">
<meta name="apple-mobile-web-app-status-bar-style" content="black-translucent">
<title>NOVA HYPERION — DEEP SPACE INTELLIGENCE TERMINAL</title>
<link href="https://fonts.googleapis.com/css2?family=Share+Tech+Mono&family=Orbitron:wght@400;600;700;800;900&display=swap" rel="stylesheet">
<script src="https://cdnjs.cloudflare.com/ajax/libs/marked/9.1.6/marked.min.js"></script>
<script>
window.MathJax = {
  tex: {
    inlineMath: [['$','$'],['\\(','\\)']],
    displayMath: [['$$','$$'],['\\[','\\]']],
    processEscapes: true
  },
  options: { ignoreHtmlClass:'no-mathjax', processHtmlClass:'mathjax-content' }
};
</script>
<script id="MathJax-script" async src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js"></script>

<style>
/* ===== LCARS PALETTE ===== */
:root {
  --bg:       #04080F;
  --gold:     #FFB300;
  --coral:    #FF6B35;
  --blue:     #7EB8FF;
  --teal:     #00D4AA;
  --pale:     #D0E8FF;
  --dimteal:  #008C6E;
  --dimblu:   #3C64A0;
  --dimgold:  #B47800;
  --purple:   #A050DC;
  --red:      #C82828;
  --indigo:   #4400FF;
  --violet:   #2A0845;
  --cobalt:   #0D1B4B;
}

* { box-sizing:border-box; margin:0; padding:0; }
html, body { background:var(--bg); font-family:'Share Tech Mono',monospace;
  overflow:hidden; height:100%; position:fixed; inset:0; }

/* ===== CANVAS STARFIELD (parallax) ===== */
#star-canvas {
  position:fixed; inset:0; width:100%; height:100%;
  pointer-events:none; z-index:0;
  display:block;
}

/* ===== APP SHELL ===== */
#app { position:fixed; inset:0; z-index:2; display:flex; flex-direction:column; }

/* ===== HEADER with LCARS styling ===== */
.hdr {
  background:linear-gradient(180deg,#08102A 0%,#04080F 100%);
  border-bottom:3px solid var(--gold); padding:10px 18px;
  flex-shrink:0; position:relative;
}
.hdr::before {
  content:''; position:absolute; bottom:-3px; left:0;
  width:220px; height:3px;
  background:linear-gradient(90deg,var(--gold),var(--coral));
}
.name {
  font-family:'Orbitron',monospace; font-size:26px; font-weight:800;
  color:var(--gold); letter-spacing:4px;
  text-shadow:0 0 12px rgba(255,179,0,.55);
}
.subtitle {
  font-size:9px; color:var(--teal); letter-spacing:2px;
  margin-top:2px;
}
.status-bar { display:flex; gap:8px; margin-top:8px; flex-wrap:wrap; align-items:center; }

/* ===== HEALTH CONDUIT (from Tkinter) ===== */
.health-conduit {
  display:flex; align-items:center; gap:4px;
  background:rgba(0,0,0,0.3); padding:2px 8px; border-radius:3px;
}
.conduit-lines {
  width:60px; height:12px; position:relative; overflow:hidden;
}
.conduit-line {
  position:absolute; height:2px; background:var(--dimteal); width:100%;
}
.conduit-line.top { top:2px; }
.conduit-line.bot { bottom:2px; }
.conduit-dash {
  position:absolute; height:2px; background:var(--teal); width:10px;
  animation:flowDash 1s linear infinite;
}
@keyframes flowDash {
  0% { left:-10px; }
  100% { left:70px; }
}
.conduit-pulse {
  width:8px; height:8px; border-radius:50%; background:var(--teal);
  animation:pulseCore 1.5s ease-in-out infinite;
}
@keyframes pulseCore {
  0%,100% { opacity:0.3; transform:scale(0.8); box-shadow:0 0 2px var(--teal); }
  50% { opacity:1; transform:scale(1.2); box-shadow:0 0 8px var(--teal); }
}
.health-processing .conduit-dash { background:var(--gold); animation-duration:0.5s; }
.health-processing .conduit-pulse { background:var(--gold); animation-duration:0.8s; }
.health-offline .conduit-line { background:var(--red); }
.health-offline .conduit-dash { display:none; }
.health-offline .conduit-pulse { background:var(--red); animation:none; opacity:0.5; }

/* ===== PILLS ===== */
.pill {
  display:flex; align-items:center; gap:6px;
  background:rgba(0,212,170,.07);
  border-left:3px solid var(--coral);
  padding:4px 11px; font-size:10px; font-weight:bold;
  color:var(--pale); white-space:nowrap;
}
.pill.ok  { border-left-color:var(--teal); }
.pill.btn { cursor:pointer; transition:background .15s; }
.pill.btn:hover  { background:rgba(0,212,170,.16); }
.pill.btn:active { transform:scale(.97); }

/* ===== PLASMA SPHERE (from Tkinter) ===== */
.plasma-sphere {
  width:16px; height:16px; position:relative; display:inline-block;
}
.plasma-core {
  width:12px; height:12px; background:#880000; border-radius:50%;
  position:absolute; top:2px; left:2px;
  box-shadow:0 0 4px var(--red);
  animation:plasmaPulse 2s ease-in-out infinite;
}
.plasma-ring {
  position:absolute; border-radius:50%; border:1px solid var(--teal);
  animation:plasmaExpand 2s ease-in-out infinite;
}
@keyframes plasmaPulse {
  0%,100% { box-shadow:0 0 2px var(--red); background:#880000; }
  50% { box-shadow:0 0 8px var(--red); background:#CC2200; }
}
@keyframes plasmaExpand {
  0% { width:12px; height:12px; top:2px; left:2px; opacity:0.8; }
  100% { width:24px; height:24px; top:-4px; left:-4px; opacity:0; }
}

/* ===== CONVERSATION ===== */
.conv {
  flex:1; overflow-y:auto; padding:16px;
  scroll-behavior:smooth; -webkit-overflow-scrolling:touch;
}
.conv::-webkit-scrollbar { width:5px; }
.conv::-webkit-scrollbar-track { background:#0A1430; }
.conv::-webkit-scrollbar-thumb { background:var(--blue); border-radius:3px; }

/* ===== MESSAGES ===== */
.msg { margin-bottom:18px; animation:fadeUp .3s ease-out forwards; opacity:0; }
@keyframes fadeUp {
  from { opacity:0; transform:translateY(14px); }
  to   { opacity:1; transform:translateY(0); }
}
.msg-hdr {
  font-size:9px; font-weight:bold; letter-spacing:1px;
  margin-bottom:4px; display:flex; align-items:center; gap:8px;
}
.msg.user      .msg-hdr { justify-content:flex-end; color:var(--gold); }
.msg.assistant .msg-hdr { color:var(--teal); }
.msg.system    .msg-hdr { justify-content:center; color:var(--coral); }

.bubble {
  max-width:80%; padding:11px 16px; border-radius:6px;
  font-size:13px; line-height:1.65; color:var(--pale); position:relative;
}
.msg.user      .bubble { background:rgba(255,107,53,.10); border-right:3px solid var(--coral); margin-left:auto; }
.msg.assistant .bubble { background:rgba(0,212,170,.05);  border-left:3px solid var(--teal); }

.save-btn {
  position:absolute; top:6px; right:6px;
  width:27px; height:27px; border-radius:4px;
  background:rgba(0,212,170,.1); border:1px solid rgba(0,212,170,.3);
  color:var(--teal); font-size:12px; cursor:pointer;
  display:flex; align-items:center; justify-content:center;
  opacity:0; transition:opacity .2s;
}
.msg.assistant:hover .save-btn { opacity:1; }

.copy-txt-btn {
  position:absolute; top:6px; left:6px;
  width:27px; height:27px; border-radius:4px;
  background:rgba(255,107,53,.1); border:1px solid rgba(255,107,53,.3);
  color:var(--coral); font-size:12px; cursor:pointer;
  display:flex; align-items:center; justify-content:center;
  opacity:0; transition:opacity .2s;
}
.msg.user:hover .copy-txt-btn { opacity:1; }

.ref-num { font-size:8px; color:var(--dimgold); margin-top:4px; text-align:right; }

/* ===== MARKDOWN ===== */
.bubble p { margin-bottom:.65em; }
.bubble p:last-child { margin-bottom:0; }
.bubble ul, .bubble ol { margin:.4em 0 .65em 1.4em; }
.bubble li { margin-bottom:.25em; }
.bubble h1,.bubble h2,.bubble h3,.bubble h4 {
  margin:1em 0 .35em; color:var(--gold);
  font-family:'Orbitron',monospace; font-weight:600;
}
.bubble strong { color:#fff; font-weight:600; }
.bubble em { color:var(--blue); }
.bubble blockquote { border-left:2px solid var(--teal); padding-left:10px; margin:.5em 0; color:var(--dimteal); }
.bubble a { color:var(--blue); text-decoration:none; }
.bubble a:hover { color:var(--teal); text-decoration:underline; }
.bubble hr { border:none; height:1px; margin:1em 0; background:linear-gradient(90deg,transparent,var(--gold),transparent); }
.bubble table { border-collapse:collapse; width:100%; margin:.65em 0; font-size:11px; }
.bubble th,.bubble td { padding:5px 9px; border:1px solid rgba(126,184,255,.2); text-align:left; }
.bubble th { background:rgba(255,179,0,.1); color:var(--gold); font-weight:600; }
.bubble code:not(pre code) {
  background:rgba(0,212,170,.12); color:var(--teal);
  padding:1px 5px; border-radius:3px; font-size:.88em; font-family:monospace;
}

/* ===== CODE BLOCKS ===== */
.code-wrap {
  margin:10px 0; border-radius:7px; overflow:hidden;
  border:1px solid rgba(0,212,170,.18); background:#060C18;
}
.code-hdr {
  display:flex; align-items:center; justify-content:space-between;
  padding:6px 12px; background:rgba(0,0,0,.35);
  border-bottom:1px solid rgba(0,212,170,.12);
}
.tl-wrap { display:flex; gap:5px; }
.tl { width:10px; height:10px; border-radius:50%; }
.tl-r{background:#FF5F57;} .tl-y{background:#FEBC2E;} .tl-g{background:#28C840;}
.cp-btn {
  opacity:0; padding:3px 10px;
  background:rgba(255,255,255,.05);
  border:1px solid rgba(126,184,255,.3); border-radius:999px;
  color:var(--blue); font-size:10px; cursor:pointer; transition:opacity .2s;
  font-family:'Share Tech Mono',monospace;
}
.code-wrap:hover .cp-btn { opacity:1; }
pre { padding:14px; overflow-x:auto; font-size:12px; margin:0; }

/* ===== WARP CORE TYPING INDICATOR ===== */
#typing-ind {
  display:none; padding:10px 18px; margin:8px;
  background:rgba(0,212,170,.04); border-left:3px solid var(--teal);
  border-radius:6px;
  align-items:center; gap:14px; font-size:10px; color:var(--teal); flex-shrink:0;
}
#typing-ind.active { display:flex; }
.warp-bars { display:flex; gap:5px; align-items:center; }
.warp-bar {
  width:5px; background:var(--gold); animation:warpPulse .8s ease-in-out infinite;
}
@keyframes warpPulse {
  0%,100% { height:6px;  opacity:.35; }
  50%      { height:20px; opacity:1; }
}
.dodec-canvas {
  width:30px; height:30px; margin-left:8px;
}
.dodec-canvas canvas {
  width:100%; height:100%;
}

/* ===== LCARS DIVIDER ===== */
.lcars-divider {
  height:40px; margin:8px 0; position:relative;
  background:linear-gradient(90deg, var(--coral) 0%, var(--coral) 30px,
              var(--amber) 30px, var(--amber) calc(100% - 30px),
              var(--coral) calc(100% - 30px), var(--coral) 100%);
  background-size:auto;
}
.lcars-divider::before {
  content:''; position:absolute; left:15px; top:12px;
  width:20px; height:16px; background:var(--coral); border-radius:4px;
}
.lcars-divider::after {
  content:''; position:absolute; right:15px; top:12px;
  width:20px; height:16px; background:var(--coral); border-radius:4px;
}
.divider-blocks {
  position:absolute; left:50%; transform:translateX(-50%);
  display:flex; gap:15px; top:8px;
}
.divider-block {
  width:18px; height:24px; background:var(--teal);
  clip-path: polygon(0% 0%, 100% 0%, 80% 100%, 20% 100%);
  animation: divCycle 3s ease-in-out infinite;
}
.divider-block:nth-child(2) { animation-delay: 1s; }
.divider-block:nth-child(3) { animation-delay: 2s; }

@keyframes divCycle {
  0%,100% { background: var(--teal); }
  50%     { background: var(--gold); }
}

/* ===== INPUT CONSOLE ===== */
.console {
  flex-shrink:0;
  background:linear-gradient(0deg,#060C18 0%,#0A1430 100%);
  border-top:3px solid var(--gold); padding:12px 18px; position:relative;
}
.console::before {
  content:''; position:absolute; top:-3px; right:0;
  width:200px; height:3px;
  background:linear-gradient(270deg,var(--gold),var(--coral));
}
.drop-zone {
  border:1px dashed rgba(126,184,255,.22); border-radius:5px;
  background:rgba(4,8,15,.7); padding:7px 14px; margin-bottom:9px;
  cursor:pointer; display:flex; align-items:center; gap:8px;
  font-size:10px; color:var(--dimteal); transition:all .25s;
}
.drop-zone.drag-over { border-color:var(--teal); background:rgba(0,212,170,.08); color:var(--teal); }
.file-prev { display:flex; flex-wrap:wrap; gap:6px; margin-bottom:8px; }
.fpi {
  background:rgba(126,184,255,.1); border-radius:4px;
  padding:3px 9px; font-size:10px; color:var(--pale);
  display:flex; align-items:center; gap:5px;
}
.fpi button { background:none; border:none; color:var(--coral); cursor:pointer; font-size:11px; padding:0 2px; }

.input-row { display:flex; gap:9px; align-items:flex-end; }
.ifield {
  flex:1; background:rgba(4,8,15,.9); border:1px solid var(--dimblu);
  border-radius:5px; padding:13px 15px; color:var(--pale);
  font-family:'Share Tech Mono',monospace; font-size:13px;
  resize:none; outline:none; transition:border-color .2s;
  min-height:48px; max-height:130px;
}
.ifield:focus { border-color:var(--blue); box-shadow:0 0 8px rgba(126,184,255,.2); }
.ifield::placeholder { color:var(--dimteal); }

.btn {
  padding:11px 18px; border:none; border-radius:5px;
  font-family:'Orbitron',monospace; font-weight:700;
  font-size:11px; cursor:pointer; letter-spacing:1px;
  transition:transform .15s; white-space:nowrap;
  -webkit-tap-highlight-color:transparent;
}
.btn:active { transform:scale(.96); }
.btn-tx  { background:linear-gradient(135deg,var(--coral),#CC4422); color:#fff; box-shadow:0 0 10px rgba(255,107,53,.35); }
.btn-hd  { background:linear-gradient(135deg,var(--purple),#6B2FA0); color:#fff; box-shadow:0 0 10px rgba(160,80,220,.35); }
.btn-mic { background:linear-gradient(135deg,var(--teal),#008C6E);   color:#fff; box-shadow:0 0 10px rgba(0,212,170,.35); }
.btn-mic.recording { background:linear-gradient(135deg,var(--red),#8C1010); animation:micPulse 1s ease-in-out infinite; }
@keyframes micPulse {
  0%,100% { box-shadow:0 0 10px rgba(200,40,40,.35); }
  50%      { box-shadow:0 0 22px rgba(200,40,40,.75); }
}
.console-foot { display:flex; justify-content:space-between; margin-top:6px; font-size:9px; color:var(--dimteal); }

/* ===== MEDIA ===== */
.media-wrap { margin:10px 0; text-align:center; }
.media-wrap img { max-width:100%; height:auto; border-radius:7px; cursor:pointer;
  box-shadow:0 2px 8px rgba(0,0,0,.4); display:block; margin:0 auto; }
.media-cap { font-size:9px; color:var(--dimteal); margin-top:4px; }

/* ===== TOAST ===== */
.toast {
  position:fixed; bottom:18px; left:50%; transform:translateX(-50%);
  padding:9px 18px; border-radius:4px; font-size:11px; font-weight:bold;
  z-index:10000; animation:slideUp .3s ease-out; white-space:nowrap;
  background:rgba(0,212,170,.95); color:var(--bg);
  font-family:'Share Tech Mono',monospace;
}
@keyframes slideUp {
  from { transform:translateX(-50%) translateY(100%); opacity:0; }
  to   { transform:translateX(-50%) translateY(0);    opacity:1; }
}

/* ===== WELCOME ===== */
.welcome { text-align:center; padding:50px 20px; animation:fadeUp .5s ease-out forwards; opacity:0; }
.delta   { font-size:70px; margin-bottom:16px; animation:dPulse 3s ease-in-out infinite; }
@keyframes dPulse { 0%,100%{opacity:.8;text-shadow:0 0 20px var(--gold);}50%{opacity:1;text-shadow:0 0 40px var(--gold);} }
.w-title { font-family:'Orbitron',monospace; font-size:30px; font-weight:800; color:var(--gold); letter-spacing:8px; margin-bottom:8px; }
.w-sub   { font-size:12px; color:var(--teal); margin-bottom:30px; }
.chip-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:10px; max-width:800px; margin:0 auto; }
.chip {
  background:rgba(126,184,255,.06); border-left:3px solid var(--blue);
  padding:9px 14px; font-size:11px; color:var(--pale); cursor:pointer;
  transition:all .2s; text-align:left; font-family:'Share Tech Mono',monospace;
}
.chip:hover { background:rgba(126,184,255,.13); transform:translateX(4px); }

@media (max-width:600px) {
  .name { font-size:18px; letter-spacing:2px; }
  .bubble { max-width:92%; font-size:12px; }
  .btn { padding:10px 12px; font-size:10px; }
  .input-row { flex-wrap:wrap; }
  .ifield { width:100%; font-size:16px; min-height:70px; }
}
</style>
</head>
<body>

<canvas id="star-canvas"></canvas>

<div id="app">

<div class="hdr">
  <div class="name">NOVA HYPERION</div>
  <div class="subtitle">DEEP SPACE INTELLIGENCE TERMINAL — NCC-74657</div>
  <div class="status-bar">
    <div class="pill ok" id="conn-pill">
      <div class="plasma-sphere" id="plasma-sphere">
        <div class="plasma-core"></div>
        <div class="plasma-ring"></div>
      </div>
      <span id="conn-lbl">ONLINE</span>
    </div>
    <div class="health-conduit" id="health-conduit">
      <div class="conduit-lines">
        <div class="conduit-line top"></div>
        <div class="conduit-line bot"></div>
        <div class="conduit-dash"></div>
      </div>
      <div class="conduit-pulse"></div>
    </div>
    <div class="pill ok"><span>⚡</span> KNOWLEDGE CORE</div>
    

    <div class="pill btn" id="tts-pill" onclick="toggleTTS()">
      <span id="tts-icon">🔊</span><span id="tts-lbl"> TTS ON</span>
    </div>
    <div class="pill btn" onclick="toggleCam()">📷 CAMERA</div>
    <div class="pill btn" onclick="clearChat()">🗑 CLEAR</div>
  </div>
</div>
<!-- PERSONALITY PICKER -->
<div style="display:flex;align-items:center;gap:6px;background:rgba(0,0,0,0.3);
  padding:2px 8px;border-radius:3px;">
  <span style="font-size:9px;color:var(--gold);letter-spacing:1px;white-space:nowrap;">
    ◈ PERSONA
  </span>
  <select id="personalitySelect" onchange="activatePersonality()"
    style="background:#1a0a00;color:#ff9900;border:1px solid #ff6600;
    border-radius:3px;padding:2px 6px;font-family:'Share Tech Mono',monospace;
    font-size:10px;outline:none;cursor:pointer;max-width:160px;">
    <option value="">— NOVA DEFAULT —</option>
  </select>
</div>


<div class="conv" id="conv">
  <div class="welcome" id="welcome">
    <div class="delta">⚛️</div>
    <div class="w-title">HYPERION ONLINE</div>
    <div class="w-sub">AWAITING STELLAR INQUIRY — ALL SYSTEMS NOMINAL</div>
    <div class="chip-grid">
      <div class="chip" onclick="useSuggestion(this)">⚡ Explain quantum entanglement</div>
      <div class="chip" onclick="useSuggestion(this)">💻 Write a Python async example</div>
      <div class="chip" onclick="useSuggestion(this)">🕳️ How do black holes form?</div>
      <div class="chip" onclick="useSuggestion(this)">⚙️ Debug my NOVA code</div>
      <div class="chip" onclick="useSuggestion(this)">💡 Creative story ideas</div>
      <div class="chip" onclick="useSuggestion(this)">📊 Analyse this dataset</div>
    </div>
  </div>
</div>

<div id="typing-ind">
  <div class="warp-bars">
    <div class="warp-bar" style="animation-delay:0s"></div>
    <div class="warp-bar" style="animation-delay:.1s"></div>
    <div class="warp-bar" style="animation-delay:.2s"></div>
    <div class="warp-bar" style="animation-delay:.3s"></div>
    <div class="warp-bar" style="animation-delay:.4s"></div>
  </div>
  <span>NEURAL SYNTHESIS IN PROGRESS...</span>
  <div class="dodec-canvas">
    <canvas id="dodec-canvas" width="30" height="30"></canvas>
  </div>
</div>

<div id="cam-panel" style="display:none;flex-shrink:0;background:#060C18;
  border-top:2px solid var(--teal);padding:8px 18px;">
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:6px;">
    <span style="font-size:10px;color:var(--teal);font-weight:bold;">📷 CAMERA PREVIEW</span>
    <button class="btn" style="padding:4px 10px;font-size:9px;background:rgba(0,212,170,.15);
      border:1px solid var(--teal);color:var(--teal);" onclick="snapAndSend()">
      📸 CAPTURE & ASK
    </button>
    <button class="btn" style="padding:4px 10px;font-size:9px;background:rgba(200,40,40,.15);
      border:1px solid var(--red);color:var(--red);margin-left:auto;" onclick="toggleCam()">
      ✕ CLOSE
    </button>
  </div>
  <video id="cam-feed" autoplay playsinline muted
    style="width:100%;max-height:200px;border-radius:6px;
    border:1px solid rgba(0,212,170,.25);background:#000;display:block;">
  </video>
</div>

<div class="console">
  <div class="lcars-divider" id="lcars-divider">
    <div class="divider-blocks">
      <div class="divider-block"></div>
      <div class="divider-block"></div>
      <div class="divider-block"></div>
    </div>
  </div>
  <div class="drop-zone" id="drop-zone">
    <span>📎</span><span>DROP FILES OR CLICK TO UPLOAD</span>
    <input type="file" id="file-input" multiple style="display:none">
  </div>
  <div class="file-prev" id="file-prev"></div>
  <div class="input-row">
    <textarea class="ifield" id="msg-in" rows="1"
      placeholder="TRANSMIT NEURAL QUERY TO HYPERION CORE..."></textarea>
    <button class="btn btn-mic" id="mic-btn" onclick="toggleMic()">🎤 VOICE</button>
    <button class="btn btn-hd" onclick="sendImagine()">✨ HOLODECK</button>
    <button class="btn btn-tx" onclick="sendMessage()">📡 TRANSMIT</button>
  </div>
  <div class="console-foot">
    <span>SHIFT+ENTER: NEW LINE · ENTER: TRANSMIT · DROP FILES TO UPLOAD</span>
    <span id="char-info"></span>
  </div>
</div>

</div>

<script>
// ============================================================
//  PARALLAX STARFIELD (Canvas-based, from Tkinter)
// ============================================================
class StarfieldCanvas {
  constructor(canvas) {
    this.canvas = canvas;
    this.ctx = canvas.getContext('2d');
    this.width = window.innerWidth;
    this.height = window.innerHeight;
    this.layers = [];
    this.shootingStars = [];
    this.mouseX = this.width / 2;
    this.mouseY = this.height / 2;
    this.offsets = [0, 0, 0];
    this.targetOffsets = [0, 0, 0];
    this.parallaxFactors = [0.008, 0.018, 0.032];
    
    this._buildStars();
    this._resize();
    window.addEventListener('resize', () => this._resize());
    document.addEventListener('mousemove', (e) => this._onMouseMove(e));
    this._animate();
    this._shootingStarLoop();
  }
  
  _buildStars() {
    const counts = [120, 80, 40];
    const colors = ['#6688AA', '#AAD4FF', '#FFFFFF'];
    const sizes = [[1,1], [1,2], [1,3]];
    
    for (let i = 0; i < 3; i++) {
      const layer = [];
      for (let j = 0; j < counts[i]; j++) {
        layer.push({
          x: Math.random() * this.width,
          y: Math.random() * this.height,
          r: sizes[i][0] + Math.random() * (sizes[i][1] - sizes[i][0]),
          color: colors[i],
          alpha: 0.3 + Math.random() * 0.7,
          twinkle: Math.random() * Math.PI * 2,
          twinkleSpeed: 0.02 + Math.random() * 0.06
        });
      }
      this.layers.push(layer);
    }
  }
  
  _resize() {
    this.width = window.innerWidth;
    this.height = window.innerHeight;
    this.canvas.width = this.width;
    this.canvas.height = this.height;
  }
  
  _onMouseMove(e) {
    this.mouseX = e.clientX;
    this.mouseY = e.clientY;
    const cx = this.width / 2;
    const cy = this.height / 2;
    const dx = (this.mouseX - cx) / cx;
    const dy = (this.mouseY - cy) / cy;
    for (let i = 0; i < 3; i++) {
      this.targetOffsets[i] = dx * this.parallaxFactors[i] * this.width;
    }
  }
  
  _spawnShootingStar() {
    const x0 = Math.random() * this.width * 0.6;
    const y0 = Math.random() * this.height * 0.4;
    const angle = 20 + Math.random() * 30;
    const rad = angle * Math.PI / 180;
    const length = 200 + Math.random() * 200;
    const dx = length * Math.cos(rad);
    const dy = length * Math.sin(rad);
    
    this.shootingStars.push({
      x0, y0, dx, dy,
      progress: 0,
      speed: 0.012 + Math.random() * 0.013,
      trail: []
    });
  }
  
  _shootingStarLoop() {
    setInterval(() => {
      if (Math.random() < 0.03) {
        this._spawnShootingStar();
      }
    }, 1000);
  }
  
  _drawShootingStar(ss) {
    const len = 8;
    for (let k = 0; k < len; k++) {
      const frac = ss.progress - k * 0.04;
      if (frac < 0 || frac > 1) continue;
      const x = ss.x0 + frac * ss.dx;
      const y = ss.y0 + frac * ss.dy;
      const alpha = Math.max(0, (1 - k * 0.18)) * (1 - ss.progress);
      const r = Math.max(0.5, 2 - k * 0.3);
      this.ctx.beginPath();
      this.ctx.arc(x, y, r, 0, Math.PI * 2);
      this.ctx.fillStyle = `rgba(204, 238, 255, ${alpha * 0.8})`;
      this.ctx.fill();
    }
    ss.progress += ss.speed;
  }
  
  _animate() {
    this.ctx.clearRect(0, 0, this.width, this.height);
    
    // Smooth parallax
    for (let i = 0; i < 3; i++) {
      this.offsets[i] += (this.targetOffsets[i] - this.offsets[i]) * 0.05;
    }
    
    // Draw stars
    for (let i = 0; i < 3; i++) {
      for (const star of this.layers[i]) {
        star.twinkle += star.twinkleSpeed;
        const brightness = 0.5 + 0.5 * Math.sin(star.twinkle);
        const x = star.x + this.offsets[i];
        const y = star.y;
        const r = star.r * (0.7 + 0.3 * brightness);
        this.ctx.beginPath();
        this.ctx.arc(x, y, r, 0, Math.PI * 2);
        this.ctx.fillStyle = star.color;
        this.ctx.globalAlpha = star.alpha * (0.5 + brightness * 0.5);
        this.ctx.fill();
      }
    }
    
    // Draw shooting stars
    this.ctx.globalAlpha = 1;
    const dead = [];
    for (const ss of this.shootingStars) {
      this._drawShootingStar(ss);
      if (ss.progress >= 1) dead.push(ss);
    }
    this.shootingStars = this.shootingStars.filter(ss => !dead.includes(ss));
    
    requestAnimationFrame(() => this._animate());
  }
}

// ============================================================
//  DODECAHEDRON ANIMATION (Warp Core)
// ============================================================
class DodecahedronAnim {
  constructor(canvas) {
    this.canvas = canvas;
    this.ctx = canvas.getContext('2d');
    this.angle = 0;
    this._animate();
  }
  
  _animate() {
    this.ctx.clearRect(0, 0, 30, 30);
    this.angle += 0.05;
    const cx = 15, cy = 15, r = 10;
    const pts = [];
    for (let k = 0; k < 6; k++) {
      const a = this.angle + k * Math.PI / 3;
      pts.push({
        x: cx + r * Math.cos(a),
        y: cy + r * 0.5 * Math.sin(a)
      });
    }
    this.ctx.strokeStyle = '#00D4AA';
    this.ctx.lineWidth = 1;
    for (let i = 0; i < 6; i++) {
      this.ctx.beginPath();
      this.ctx.moveTo(pts[i].x, pts[i].y);
      this.ctx.lineTo(pts[(i+1)%6].x, pts[(i+1)%6].y);
      this.ctx.stroke();
    }
    this.ctx.beginPath();
    this.ctx.moveTo(pts[0].x, pts[0].y);
    this.ctx.lineTo(pts[3].x, pts[3].y);
    this.ctx.stroke();
    this.ctx.beginPath();
    this.ctx.moveTo(pts[1].x, pts[1].y);
    this.ctx.lineTo(pts[4].x, pts[4].y);
    this.ctx.stroke();
    this.ctx.beginPath();
    this.ctx.moveTo(pts[2].x, pts[2].y);
    this.ctx.lineTo(pts[5].x, pts[5].y);
    this.ctx.stroke();
    
    requestAnimationFrame(() => this._animate());
  }
}

// ============================================================
//  AMBIENT HUM (Web Audio)
// ============================================================
let ambientContext = null;
let ambientSource = null;
let ambientEnabled = true;
let ambientGain = null;

async function initAmbient() {
  if (ambientContext) return;
  ambientContext = new (window.AudioContext || window.webkitAudioContext)();
  ambientGain = ambientContext.createGain();
  ambientGain.gain.value = 0.03;
  ambientGain.connect(ambientContext.destination);
  
  // Create 40Hz base + 400Hz harmonic + 80Hz sub
  const sampleRate = ambientContext.sampleRate;
  const bufferSize = sampleRate * 2;
  const buffer = ambientContext.createBuffer(2, bufferSize, sampleRate);
  
  for (let ch = 0; ch < 2; ch++) {
    const data = buffer.getChannelData(ch);
    for (let i = 0; i < bufferSize; i++) {
      const t = i / sampleRate;
      const base = 0.018 * Math.sin(2 * Math.PI * 40 * t);
      const harm = 0.012 * Math.sin(2 * Math.PI * 400 * t);
      const sub = 0.008 * Math.sin(2 * Math.PI * 80 * t);
      const mod = 1 + 0.05 * Math.sin(2 * Math.PI * 0.25 * t);
      data[i] = (base + harm + sub) * mod;
    }
  }
  
  const playHum = () => {
    if (!ambientEnabled || !ambientContext) return;
    const src = ambientContext.createBufferSource();
    src.buffer = buffer;
    src.loop = true;
    src.connect(ambientGain);
    src.start();
    ambientSource = src;
    src.onended = () => { if (ambientEnabled) playHum(); };
  };
  
  playHum();
  console.log('Ambient hum initialized');
}

function toggleAmbient() {
  ambientEnabled = !ambientEnabled;
  const icon = document.getElementById('ambient-icon');
  const lbl = document.getElementById('ambient-lbl');
  if (ambientEnabled) {
    icon.textContent = '🌊';
    lbl.textContent = ' HUM ON';
    if (ambientContext) {
      ambientContext.resume();
      initAmbient();
    } else {
      initAmbient();
    }
  } else {
    icon.textContent = '🔇';
    lbl.textContent = ' HUM OFF';
    if (ambientSource) {
      try { ambientSource.stop(); } catch(e) {}
    }
    if (ambientContext) ambientContext.suspend();
  }
  fetch('/api/ambient', { method: 'POST' });
}

// ============================================================
//  HEALTH CONDUIT STATUS
// ============================================================
function setHealthStatus(status) {
  const conduit = document.getElementById('health-conduit');
  conduit.classList.remove('health-processing', 'health-offline');
  if (status === 'processing') conduit.classList.add('health-processing');
  else if (status === 'offline') conduit.classList.add('health-offline');
}

// ============================================================
//  EXISTING FUNCTIONS (preserved from original)
// ============================================================
const conv = document.getElementById('conv');
const msgIn = document.getElementById('msg-in');
let lastCount = 0, lastContent = '', isThinking = false, userScrolled = false;
let audioCtx = null, curAudio = null;
let mediaRec = null, audioChunks = [], isRec = false;
let pendingFiles = [], refCounter = 47291;
let camStream = null;

function toast(msg, err=false) {
  const t = document.createElement('div');
  t.className = 'toast';
  t.style.background = err ? 'rgba(200,40,40,.95)' : 'rgba(0,212,170,.95)';
  t.style.color = err ? '#fff' : '#04080F';
  t.textContent = msg;
  document.body.appendChild(t);
  setTimeout(() => t.remove(), 3000);
}

async function checkHealth() {
  try {
    const r = await fetch('/api/ping');
    const lbl = document.getElementById('conn-lbl');
    if (r.ok) { lbl.textContent = 'ONLINE'; }
    else { lbl.textContent = 'DEGRADED'; setHealthStatus('offline'); }
  } catch {
    document.getElementById('conn-lbl').textContent = 'OFFLINE';
    setHealthStatus('offline');
  }
}
checkHealth();
setInterval(checkHealth, 4000);

function unlockAudio() {
  if (audioCtx) return;
  audioCtx = new (window.AudioContext || window.webkitAudioContext)();
  const buf = audioCtx.createBuffer(1, 1, 22050);
  const src = audioCtx.createBufferSource();
  src.buffer = buf;
  src.connect(audioCtx.destination);
  src.start(0);
  audioCtx.resume();
}
document.addEventListener('touchstart', unlockAudio, {once:true});
document.addEventListener('click', unlockAudio, {once:true});

function stopAudio() {
  if (curAudio) { try{curAudio.stop();}catch(e){} curAudio = null; }
}

async function playTTS(text) {
  stopAudio();
  try {
    const r = await fetch('/api/speak', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({text})
    });
    if (!r.ok || r.status === 204) return;
    const buf = await r.arrayBuffer();
    if (!buf.byteLength) return;
    if (!audioCtx || audioCtx.state === 'closed')
      audioCtx = new (window.AudioContext || window.webkitAudioContext)();
    await audioCtx.resume();
    const decoded = await audioCtx.decodeAudioData(buf);
    const src = audioCtx.createBufferSource();
    curAudio = src;
    src.buffer = decoded;
    src.connect(audioCtx.destination);
    src.start(0);
    src.onended = () => { curAudio = null; };
  } catch(e) { console.warn('TTS:', e); curAudio = null; }
}

async function toggleTTS() {
  const r = await fetch('/api/tts', {method:'POST'});
  const d = await r.json();
  const on = d.tts;
  document.getElementById('tts-icon').textContent = on ? '🔊' : '🔇';
  document.getElementById('tts-lbl').textContent = on ? ' TTS ON' : ' TTS OFF';
  document.getElementById('tts-pill').style.borderLeftColor = on ? 'var(--teal)' : 'var(--coral)';
  if (!on) stopAudio();
  toast(on ? '🔊 TTS enabled' : '🔇 TTS muted');
}

async function toggleCam() {
  const panel = document.getElementById('cam-panel');
  const feed = document.getElementById('cam-feed');
  if (camStream) {
    camStream.getTracks().forEach(t => t.stop());
    camStream = null;
    feed.srcObject = null;
    panel.style.display = 'none';
    return;
  }
  try {
    camStream = await navigator.mediaDevices.getUserMedia({
      video: { facingMode: 'environment', width: {ideal:1280}, height: {ideal:720} },
      audio: false
    });
    feed.srcObject = camStream;
    panel.style.display = 'block';
    toast('📷 Camera preview live');
  } catch(e) {
    toast('❌ Camera access denied', true);
  }
}

async function snapAndSend() {
  if (!camStream) { toast('❌ Camera not active', true); return; }
  const feed = document.getElementById('cam-feed');
  const canvas = document.createElement('canvas');
  canvas.width = feed.videoWidth || 640;
  canvas.height = feed.videoHeight || 480;
  canvas.getContext('2d').drawImage(feed, 0, 0);
  
  camStream.getTracks().forEach(t => t.stop());
  camStream = null;
  feed.srcObject = null;
  document.getElementById('cam-panel').style.display = 'none';
  
  canvas.toBlob(async blob => {
    if (!blob) { toast('❌ Capture failed', true); return; }
    const fd = new FormData();
    fd.append('files', blob, 'camera_snapshot.jpg');
    toast('⏳ Uploading...');
    try {
      const r = await fetch('/api/upload', { method:'POST', body:fd });
      const d = await r.json();
      if (d.status === 'success') {
        const f = d.files[0];
        pendingFiles.push(f);
        showFilePreviews([f]);
        const prompt = msgIn.value.trim() || 'Describe everything you can see in this image.';
        msgIn.value = prompt;
        toast('✓ Snapshot ready');
        setTimeout(() => sendMessage(), 300);
      }
    } catch(e) { toast('❌ Upload failed', true); }
  }, 'image/jpeg', 0.92);
}

function toggleMic() {
  if (isRec) { stopMic(); return; }
  navigator.mediaDevices.getUserMedia({audio:true})
    .then(stream => {
      audioChunks = [];
      mediaRec = new MediaRecorder(stream);
      mediaRec.ondataavailable = e => { if (e.data.size > 0) audioChunks.push(e.data); };
      mediaRec.onstop = async () => {
        stream.getTracks().forEach(t => t.stop());
        await sendVoice(new Blob(audioChunks, {type:'audio/webm'}));
      };
      mediaRec.start();
      isRec = true;
      const b = document.getElementById('mic-btn');
      b.classList.add('recording');
      b.textContent = '⏹ STOP';
      toast('🎤 Recording...');
    })
    .catch(() => toast('❌ Microphone access denied', true));
}

function stopMic() {
  if (mediaRec && mediaRec.state !== 'inactive') mediaRec.stop();
  isRec = false;
  const b = document.getElementById('mic-btn');
  b.classList.remove('recording');
  b.textContent = '🎤 VOICE';
}

async function sendVoice(blob) {
  toast('⏳ Transcribing...');
  try {
    const r = await fetch('/api/voice', {
      method:'POST',
      headers:{'Content-Type':'audio/webm', 'Content-Length':blob.size},
      body:blob
    });
    const d = await r.json();
    if (d.transcript) {
      msgIn.value = d.transcript;
      msgIn.dispatchEvent(new Event('input'));
      toast('✓ Transcript ready');
    } else toast('⚠ No speech detected', true);
  } catch { toast('❌ Voice error', true); }
}

function addCodeCopy(container) {
  container.querySelectorAll('pre').forEach(pre => {
    if (pre.closest('.code-wrap')) return;
    const wrap = document.createElement('div');
    wrap.className = 'code-wrap';
    pre.parentNode.insertBefore(wrap, pre);
    wrap.appendChild(pre);
    const hdr = document.createElement('div');
    hdr.className = 'code-hdr';
    hdr.innerHTML = '<div class="tl-wrap"><div class="tl tl-r"></div><div class="tl tl-y"></div><div class="tl tl-g"></div></div>';
    const cb = document.createElement('button');
    cb.className = 'cp-btn';
    cb.textContent = '📋 COPY';
    cb.onclick = e => {
      e.stopPropagation();
      navigator.clipboard.writeText(pre.textContent).then(() => {
        cb.textContent = '✓ COPIED';
        setTimeout(() => cb.textContent = '📋 COPY', 2000);
        toast('✓ Code copied');
      });
    };
    hdr.appendChild(cb);
    wrap.insertBefore(hdr, pre);
  });
}

async function renderMath(el) {
  if (!window.MathJax || !el) return;
  try { MathJax.typesetClear([el]); await MathJax.typesetPromise([el]); }
  catch(e) { console.warn('MathJax:', e); }
}


// ═══════════════════════════════════════════════════════════
//  VIDEO PLAYER
// ═══════════════════════════════════════════════════════════
function buildVideo(url, name) {
  const wrap = document.createElement('div');
  wrap.style.cssText = 'background:#060C18;border-radius:6px;overflow:hidden;margin-top:8px;border:1px solid rgba(0,212,170,.2);position:relative;';

  const vid = document.createElement('video');
  vid.style.cssText = 'width:100%;max-height:340px;display:block;background:#000;';
  vid.muted = true;
  vid.autoplay = false;
  vid.playsInline = true;
  vid.controls = false;

  // ── MIME type from extension ──────────────────────────────
  const source = document.createElement('source');
  source.src = url;
  const ext = name.split('.').pop().toLowerCase();
  const mimeMap = {'mp4':'video/mp4','m4v':'video/mp4','mov':'video/mp4',
                   'mkv':'video/webm','webm':'video/webm','avi':'video/x-msvideo'};
  source.type = mimeMap[ext] || 'video/mp4';
  vid.appendChild(source);
  vid.addEventListener('loadedmetadata', () => {
    if(vid.videoWidth === 0) {
        vid.dispatchEvent(new Event('error'));
    }
});
  wrap.appendChild(vid);

  // ── UNMUTE button ─────────────────────────────────────────
  const soundBtn = document.createElement('button');
  soundBtn.innerHTML = '🔇 UNMUTE';
  soundBtn.style.cssText = `
    position:absolute; top:10px; right:10px;
    background:rgba(0,0,0,0.8); color:var(--teal);
    border:1px solid var(--teal); border-radius:4px;
    padding:6px 12px; font-size:10px; cursor:pointer;
    font-family:'Share Tech Mono',monospace;
    font-weight:bold; z-index:20; backdrop-filter:blur(4px); transition:all 0.2s;
  `;
  soundBtn.onmouseenter = () => soundBtn.style.background = 'rgba(0,212,170,0.3)';
  soundBtn.onmouseleave = () => soundBtn.style.background = 'rgba(0,0,0,0.8)';
  let isMuted = true;
  soundBtn.onclick = (e) => {
    e.stopPropagation(); isMuted = !isMuted; vid.muted = isMuted;
    soundBtn.innerHTML = isMuted ? '🔇 UNMUTE' : '🔊 MUTE';
    if(typeof toast==='function') toast(isMuted ? 'Sound muted' : 'Sound enabled');
  };
  wrap.appendChild(soundBtn);

  // ── Play overlay ──────────────────────────────────────────
  const playOverlay = document.createElement('div');
  playOverlay.innerHTML = '▶';
  playOverlay.style.cssText = `
    position:absolute; top:50%; left:50%; transform:translate(-50%,-50%);
    width:50px; height:50px; background:rgba(0,212,170,0.85);
    border-radius:50%; display:flex; align-items:center; justify-content:center;
    font-size:24px; color:#04080F; cursor:pointer;
    opacity:0.8; transition:all 0.2s; z-index:15;
  `;
  playOverlay.onmouseenter = () => playOverlay.style.opacity = '1';
  playOverlay.onmouseleave = () => playOverlay.style.opacity = '0.8';
  playOverlay.onclick = (e) => {
    e.stopPropagation();
    if(vid.paused){ vid.play().catch(err=>console.error('Play error:',err)); }
  };
  vid.onclick = () => { if(vid.paused) vid.play().catch(e=>{}); else vid.pause(); };
  wrap.appendChild(playOverlay);

  // ── Controls bar ──────────────────────────────────────────
  const ctrl = document.createElement('div');
  ctrl.style.cssText = 'padding:7px 10px;background:rgba(0,0,0,.5);';

  const pw = document.createElement('div');
  pw.style.cssText = 'height:4px;background:rgba(255,255,255,.15);border-radius:2px;margin-bottom:7px;cursor:pointer;';
  const pb = document.createElement('div');
  pb.style.cssText = 'height:100%;width:0%;background:linear-gradient(90deg,var(--gold),var(--coral));border-radius:2px;pointer-events:none;';
  pw.appendChild(pb);
  pw.onclick = e => { const r=pw.getBoundingClientRect(); vid.currentTime=((e.clientX-r.left)/r.width)*vid.duration; };
  ctrl.appendChild(pw);

  const row = document.createElement('div');
  row.style.cssText = 'display:flex;gap:6px;align-items:center;flex-wrap:wrap;';
  const bs = 'border-radius:3px;padding:3px 8px;cursor:pointer;font-size:10px;border:1px solid rgba(0,212,170,.25);background:rgba(0,212,170,.08);color:var(--pale);font-family:"Share Tech Mono",monospace;';

  const playPauseBtn = document.createElement('button');
  playPauseBtn.style.cssText = bs;
  playPauseBtn.textContent = '▶ PLAY';
  playPauseBtn.onclick = () => {
    if(vid.paused){ vid.play().catch(e=>{}); playPauseBtn.textContent='⏸ PAUSE'; }
    else { vid.pause(); playPauseBtn.textContent='▶ PLAY'; }
  };

  vid.onplay = () => {
    playPauseBtn.textContent = '⏸ PAUSE';
    playOverlay.style.opacity = '0'; playOverlay.style.pointerEvents = 'none';
  };
  vid.onpause = () => {
    playPauseBtn.textContent = '▶ PLAY';
    playOverlay.style.opacity = '0.8'; playOverlay.style.pointerEvents = 'auto';
  };
  row.appendChild(playPauseBtn);

  [[-10,'⏪ -10s'],[10,'+10s ⏩']].forEach(([d,l])=>{
    const b=document.createElement('button'); b.style.cssText=bs; b.textContent=l;
    b.onclick=()=>{ vid.currentTime=Math.min(vid.duration,Math.max(0,vid.currentTime+d)); };
    row.appendChild(b);
  });

  const tim = document.createElement('span');
  tim.style.cssText = 'color:var(--dimteal);font-size:10px;min-width:80px;';
  tim.textContent = '0:00/0:00';
  row.appendChild(tim);

  const spd = document.createElement('span');
  spd.style.cssText = 'color:var(--dimteal);font-size:10px;'; spd.textContent = 'SPD:';
  row.appendChild(spd);
  const spdBtns = [];
  [0.5,1,1.5,2].forEach(s=>{
    const b=document.createElement('button'); b.style.cssText=bs; b.textContent=s+'x';
    b.onclick=()=>{ vid.playbackRate=s; spdBtns.forEach(x=>x.style.background='rgba(0,212,170,.08)'); b.style.background='rgba(0,212,170,.3)'; };
    if(s===1) b.style.background='rgba(0,212,170,.3)';
    spdBtns.push(b); row.appendChild(b);
  });

  const fsB = document.createElement('button');
  fsB.style.cssText = bs+'margin-left:auto;'; fsB.textContent = '⛶ FULL';
  fsB.onclick = () => { if(vid.requestFullscreen) vid.requestFullscreen(); };
  row.appendChild(fsB);

  const fmt = s => Math.floor(s/60)+':'+String(Math.floor(s%60)).padStart(2,'0');
  vid.addEventListener('timeupdate',()=>{
    if(vid.duration){ pb.style.width=(vid.currentTime/vid.duration*100)+'%'; tim.textContent=fmt(vid.currentTime)+'/'+fmt(vid.duration); }
  });

  ctrl.appendChild(row);

  // ── Caption + VLC fallback ────────────────────────────────
  const cap = document.createElement('div');
  cap.className = 'media-cap'; cap.textContent = '🎬 ' + name;
  ctrl.appendChild(cap);
  wrap.appendChild(ctrl);

  // ── Error handler — show VLC button if browser can't play ─
  vid.addEventListener('error', () => {
    const errDiv = document.createElement('div');
    errDiv.style.cssText = 'padding:16px;text-align:center;';
    errDiv.innerHTML = `
      <div style="color:#ff6b6b;margin-bottom:10px;">
        ⚠️ Browser cannot play this codec (likely MPEG-4 / Xvid).<br>
        <span style="font-size:9px;color:var(--dimteal);">Only H.264 MP4 plays in browsers.</span>
      </div>
    `;
    const vlcBtn = document.createElement('button');
    vlcBtn.style.cssText = bs + 'font-size:12px;padding:8px 16px;';
    vlcBtn.textContent = '🎬 Open in VLC / Desktop Player';
    vlcBtn.onclick = () => {
      fetch('/api/send', {method:'POST',
        headers:{'Content-Type':'application/json'},
        body:JSON.stringify({message:'open ' + name})});
      if(typeof toast==='function') toast('Opening in desktop player...');
    };
    errDiv.appendChild(vlcBtn);

    // Also offer download link
    const dl = document.createElement('a');
    dl.href = url; dl.download = name;
    dl.style.cssText = 'display:block;margin-top:8px;font-size:9px;color:var(--dimteal);';
    dl.textContent = '⬇️ Download file';
    errDiv.appendChild(dl);

    wrap.insertBefore(errDiv, vid);
    vid.style.display = 'none';
    soundBtn.style.display = 'none';
    playOverlay.style.display = 'none';
  });

  return wrap;
}



// ═══════════════════════════════════════════════════════════
//  SAVE RESPONSE
// ═══════════════════════════════════════════════════════════
const SAVECSS=`
*{box-sizing:border-box;}
body{background:#04080F;color:#D0E8FF;font-family:'Share Tech Mono',monospace;padding:40px 20px;line-height:1.75;}
.c{max-width:960px;margin:0 auto;background:rgba(10,20,40,.75);border:1px solid rgba(0,212,170,.25);border-radius:10px;padding:32px;}
h1,h2,h3,h4{font-family:'Orbitron',monospace;color:#FFB300;margin:1em 0 .4em;}
h1{font-size:1.5em;} h2{font-size:1.2em;} h3{font-size:1em;}
h1:first-child,h2:first-child{margin-top:0;}
.ts{color:#3C64A0;font-size:10px;margin-bottom:24px;}
p{margin-bottom:.75em;}
hr{border:none;height:1px;background:linear-gradient(90deg,transparent,#FFB300,transparent);margin:1.2em 0;}
pre{background:#060C18;padding:14px;border-radius:6px;overflow-x:auto;border:1px solid rgba(0,212,170,.2);margin:.75em 0;}
code{background:rgba(0,212,170,.12);color:#00D4AA;padding:1px 5px;border-radius:3px;font-size:.9em;}
blockquote{border-left:3px solid #00D4AA;padding-left:12px;margin:.75em 0;color:#008C6E;}
a{color:#7EB8FF;}
ul,ol{margin:.5em 0 .75em 1.4em;}
li{margin-bottom:.25em;}
strong{color:#fff;}
table{border-collapse:collapse;width:100%;margin:.85em 0;font-size:13px;overflow-x:auto;display:block;}
thead tr{background:rgba(255,179,0,.12);}
th{padding:7px 11px;border:1px solid rgba(126,184,255,.25);color:#FFB300;font-weight:bold;text-align:left;}
td{padding:6px 11px;border:1px solid rgba(126,184,255,.18);text-align:left;}
tbody tr:nth-child(even){background:rgba(255,255,255,.03);}
img{max-width:100%;border-radius:6px;}
.pb{position:fixed;top:14px;right:14px;padding:8px 18px;
background:linear-gradient(135deg,#FF6B35,#CC4422);border:none;border-radius:4px;
color:#fff;cursor:pointer;font-family:'Orbitron',monospace;font-weight:bold;font-size:11px;z-index:999;}
@media print{.pb{display:none;}body{background:#fff;color:#000;}.c{background:#fff;border:none;}}
`;

function saveResponse(html){
  const action=confirm('Save response?\n\nOK = download HTML file\nCancel = open new tab');
  const full=`<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8">
<title>Nova Response — ${new Date().toLocaleString()}</title>
<link href="https://fonts.googleapis.com/css2?family=Orbitron:wght@700&family=Share+Tech+Mono&display=swap" rel="stylesheet">
<script>window.MathJax={tex:{inlineMath:[['$','$'],['\\\\(','\\\\)']],displayMath:[['$$','$$'],['\\\\[','\\\\]']],processEscapes:true},options:{skipHtmlTags:['script','noscript','style','textarea','pre']}};<\/script>
<script src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js"><\/script>
<style>${SAVECSS}</style></head><body>
<button class="pb" onclick="window.print()">🖨 PRINT / PDF</button>
<div class="c">
<h1>⚡ NOVA — KNOWLEDGE COMPUTER</h1>
<div class="ts">Generated: ${new Date().toLocaleString()}</div>
${html}
</div></body></html>`;
  if(action){
    const a=document.createElement('a');
    a.href=URL.createObjectURL(new Blob([full],{type:'text/html'}));
    a.download=`nova_${Date.now()}.html`; a.click(); URL.revokeObjectURL(a.href);
    toast('✓ Response saved');
  } else {
    const w=window.open(); w.document.write(full); w.document.close();
    toast('✓ Opened in new tab');
  }
}


// ═══════════════════════════════════════════════════════════
//  ADD MESSAGE
// ═══════════════════════════════════════════════════════════
function addMessage(role, content, scroll=true){
  const welcome=document.getElementById('welcome');
  if(welcome && role!=='system') welcome.remove();

  refCounter++;
  const ref='REF-'+refCounter+'-'+['ALPHA','BETA','GAMMA'][refCounter%3];
  const time=new Date().toLocaleTimeString('en-US',{hour:'2-digit',minute:'2-digit'});

  const row=document.createElement('div'); row.className='msg '+role;
  const label={user:'HUMAN',assistant:'COMPUTER',system:'SYSTEM'}[role]||role.toUpperCase();

  row.innerHTML=
    `<div class="msg-hdr"><span>${label}</span><span style="opacity:.55;font-size:8px">${time}</span></div>`+
    `<div class="bubble"></div>`+
    (role==='assistant'?`<div class="ref-num">${ref}</div>`:'');

  const bubble=row.querySelector('.bubble');

  // Save button (assistant)
  // We re-render from the raw `content` string (with $...$ intact) rather than
  // cloning the live DOM, so the saved file's MathJax can render equations fresh.
  if(role==='assistant'){
    const sb=document.createElement('button'); sb.className='save-btn'; sb.title='Save response'; sb.textContent='💾';
    sb.onclick=()=>{
      let saveHtml;
      if(window.marked){
        const maths=[];
        let s=content;
        const ph=m=>{ maths.push(m); return '@@M'+(maths.length-1)+'@@'; };
        s=s.replace(/\$\$[\s\S]+?\$\$/g,ph);
        s=s.replace(/\$[^\$\n]+?\$/g,ph);
        // Prevent setext headings from ============ separator lines
        s = s.replace(/^={3,}$/gm, '---');
        let html=marked.parse(s);
        html=html.replace(/@@M(\d+)@@/g,(_,i)=>maths[+i]);
        html=html.replace(/<code>(https?:\/\/[^<]+)<\/code>/g,'<a href="$1" target="_blank">$1</a>');
      saveHtml=html;
      } else {
        saveHtml=content.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/\n/g,'<br>');
      }
      saveResponse(saveHtml);
    };
    bubble.appendChild(sb);
  }

  // Copy-text button (user)
  if(role==='user'){
    const cb=document.createElement('button'); cb.className='copy-txt-btn'; cb.title='Copy text'; cb.textContent='📋';
    cb.onclick=()=>{ navigator.clipboard.writeText(content).then(()=>toast('✓ Copied')); };
    bubble.appendChild(cb);
  }

  const trim=(content||'').trim();

  // ── special message types ────────────────────────────────
if(trim.startsWith('[IMAGE:')){
    const firstLine=trim.split('\n')[0];
    const fn=firstLine.slice(7,-1);
    const w=document.createElement('div'); w.className='media-wrap';
    const img=document.createElement('img'); img.src='/images/'+fn;
    img.onclick=()=>window.open('/images/'+fn);
    img.onload=()=>{ if(!userScrolled) conv.scrollTop=conv.scrollHeight; };
    const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='📷 Click to enlarge';
    w.appendChild(img); w.appendChild(cap); bubble.appendChild(w);
    // Render any text after the image tag
    const remainder=trim.slice(firstLine.length).trim();
    if(remainder){
        const contentDiv=document.createElement('div');
        if(window.marked){
            const maths=[];
            let s=remainder;
            const ph=m=>{ maths.push(m); return '@@M'+(maths.length-1)+'@@'; };
            s=s.replace(/\$\$[\s\S]+?\$\$/g,ph);
            s=s.replace(/\$[^\$\n]+?\$/g,ph);
            let html=marked.parse(s);
            html=html.replace(/@@M(\d+)@@/g,(_,i)=>maths[+i]);
            html=html.replace(/<code>(https?:\/\/[^<]+)<\/code>/g,'<a href="$1" target="_blank">$1</a>');
            contentDiv.innerHTML = html;
        } else {
            contentDiv.textContent=remainder;
        }
        bubble.appendChild(contentDiv);
    }

} else if(trim.startsWith('[DIAGRAM:')){
    const firstLine=trim.split('\n')[0];
    const fn=firstLine.slice(9,-1);
    
    
    const w=document.createElement('div'); w.className='media-wrap';
    const img=document.createElement('img'); img.src='/images/'+fn;
    img.onclick=()=>window.open('/images/'+fn);
    img.onload=()=>{ if(!userScrolled) conv.scrollTop=conv.scrollHeight; };
    const cap=document.createElement('div'); cap.className='media-cap'; cap.textContent='📊 Generated Diagram';
    w.appendChild(img); w.appendChild(cap); bubble.appendChild(w);

} else if(content && content.includes('[AUDIO:')){
    const m=content.match(/\[AUDIO:(.+?)\]/);
    if(m){
      const url='/api/stream?file='+encodeURIComponent(m[1]);
      const name=m[1].split(/[/\\]/).pop();
      const ext=name.split('.').pop().toLowerCase();
      const mimeMap={'mp3':'audio/mpeg','wav':'audio/wav','ogg':'audio/ogg',
                     'm4a':'audio/mp4','flac':'audio/flac','wma':'audio/x-ms-wma'};
      const mime=mimeMap[ext]||'audio/mpeg';

      const txt=document.createElement('div');
      txt.textContent=content.replace(/\[AUDIO:.+?\]/,'').trim();
      bubble.appendChild(txt);

      const w=document.createElement('div'); w.className='media-wrap';
      const aud=document.createElement('audio');
      aud.controls=true; aud.style.width='100%'; aud.preload='metadata';
      const src=document.createElement('source');
      src.src=url; src.type=mime;
      aud.appendChild(src);
      const cap=document.createElement('div');
      cap.className='media-cap'; cap.textContent='🎵 '+name;
      w.appendChild(aud); w.appendChild(cap); bubble.appendChild(w);
    }
  } else if(content && content.includes('[VIDEO:')){
    const m=content.match(/\[VIDEO:(.+?)\]/);
    if(m){
      const url='/api/stream?file='+encodeURIComponent(m[1]);
      const name=m[1].split(/[/\\]/).pop();
      const txt=document.createElement('div'); txt.textContent=content.replace(/\[VIDEO:.+?\]/,'').trim();
      bubble.appendChild(txt); bubble.appendChild(buildVideo(url,name));
    }
  } else if(content && content.includes('[VIDEO:')){
    const m=content.match(/\[VIDEO:(.+?)\]/);
    if(m){
      const url='/api/stream?file='+encodeURIComponent(m[1]);
      const name=m[1].split(/[/\\]/).pop();
      const txt=document.createElement('div'); txt.textContent=content.replace(/\[VIDEO:.+?\]/,'').trim();
      bubble.appendChild(txt); bubble.appendChild(buildVideo(url,name));
    }

  } else if(content && content.includes('[PLOT:')){
    const m=content.match(/\[PLOT:(.+?)\]/);
    if(m){
      const txt=document.createElement('div');
      txt.textContent=content.replace(/\[PLOT:.+?\]/,'').trim();
      bubble.appendChild(txt);
      const w=document.createElement('div'); w.className='media-wrap';
      w.style.cssText='text-align:left;';
      const iframe=document.createElement('iframe');
      iframe.src='/plots/'+m[1];
      iframe.style.cssText='width:100%;height:500px;border:none;border-radius:6px;background:#060C18;';
      iframe.setAttribute('loading','lazy');
      const cap=document.createElement('div');
      cap.className='media-cap'; cap.textContent='📊 Interactive Plot — scroll/zoom to explore';
      w.appendChild(iframe); w.appendChild(cap); bubble.appendChild(w);
    }
  } else {
    // Markdown + Math — always append into a child div so the save-btn
    // DOM node (with its closure-based onclick) is never serialised/destroyed.
    const contentDiv = document.createElement('div');
    if(window.marked && role==='assistant'){
      const maths=[];
      let s=content;
      const ph=m=>{ maths.push(m); return '@@M'+(maths.length-1)+'@@'; };
      s=s.replace(/\$\$[\s\S]+?\$\$/g,ph);
      s=s.replace(/\$[^\$\n]+?\$/g,ph);
      let html=marked.parse(s);
      html=html.replace(/@@M(\d+)@@/g,(_,i)=>maths[+i]);
      contentDiv.innerHTML = html;
    } else {
      contentDiv.innerHTML = content
        .replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')
        .replace(/\n/g,'<br>');
    }
    bubble.appendChild(contentDiv);
    addCodeCopy(bubble);
    if(window.MathJax && role==='assistant' &&
       (contentDiv.innerHTML.includes('$') || contentDiv.innerHTML.includes('\\(')))
      setTimeout(()=>renderMath(bubble),80);
  }

  conv.appendChild(row);
  if(scroll && !userScrolled) row.scrollIntoView({behavior:'smooth',block:'end'});
}


function setThinking(on) {
  const ind = document.getElementById('typing-ind');
  ind.classList.toggle('active', on);
  const lbl = document.getElementById('conn-lbl');
  if (on) { lbl.textContent = 'PROCESSING...'; setHealthStatus('processing'); }
  else { lbl.textContent = 'ONLINE'; setHealthStatus('online'); }
}

function startPoll() {
  let wait = 0, sawThinking = false, finished = false;
  const iv = setInterval(async () => {
    if (finished) return;
    wait++;
    try {
      const [hR, sR] = await Promise.all([fetch('/api/history'), fetch('/api/state')]);
      const hist = await hR.json();
      const state = await sR.json();
      if (state.thinking) sawThinking = true;
      if (hist.length > lastCount) {
        finished = true;
        clearInterval(iv);
        setThinking(false);
        isThinking = false;
        loadHistory();
      } else if (sawThinking && !state.thinking && wait > 4) {
        finished = true;
        clearInterval(iv);
        setTimeout(() => { setThinking(false); isThinking = false; loadHistory(); }, 2000);
      } else if (wait > 1200) {
        finished = true;
        clearInterval(iv);
        setThinking(false);
        isThinking = false;
      }
    } catch(e) {}
  }, 500);
}

async function sendMessage() {
  if (isThinking) return;
  const raw = msgIn.value.trim();
  const payload = buildFilePayload();
  if (!raw && !payload) return;
  const full = payload ? (raw ? payload + 'User comment: ' + raw : payload + 'Please analyse these files.') : raw;
  
  addMessage('user', raw || (payload ? '📎 (files attached)' : ''));
  msgIn.value = '';
  msgIn.style.height = 'auto';
  document.getElementById('char-info').textContent = '';
  isThinking = true;
  setThinking(true);
  lastCount++;
  
  try {
    await fetch('/api/send', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({message: full})
    });
    startPoll();
  } catch {
    setThinking(false);
    isThinking = false;
    addMessage('system', 'COMMS FAILURE — CHECK SERVER');
  }
}

async function sendImagine() {
  if (isThinking) return;
  const raw = msgIn.value.trim();
  const payload = buildFilePayload();
  if (!raw && !payload) return;
  const base = payload ? (raw ? payload + 'User comment: ' + raw : payload + 'Analyse these.') : raw;
  
  addMessage('user', '✨ [HOLODECK] ' + (raw || '(files)'));
  msgIn.value = '';
  msgIn.style.height = 'auto';
  isThinking = true;
  setThinking(true);
  lastCount++;
  
  try {
    await fetch('/api/imagine', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({message: base})
    });
    startPoll();
  } catch {
    setThinking(false);
    isThinking = false;
    addMessage('system', 'HOLODECK OFFLINE');
  }
}

function useSuggestion(chip) {
  msgIn.value = chip.textContent.trim();
  sendMessage();
}

async function clearChat() {
  if (!confirm('CLEAR ALL COMPUTER LOGS?')) return;
  await fetch('/api/clear', {method:'POST'});
  conv.innerHTML = '<div class="welcome"><div class="delta">🔬</div><div class="w-title">HYPERION CORE ONLINE</div><div class="w-sub">LOGS CLEARED — AWAITING INQUIRY</div></div>';
  lastCount = 0;
  lastContent = '';
  toast('✓ Logs cleared');
}
// ── Personality picker ──────────────────────────────────────
let _personalities = [];

async function loadPersonalities() {
  try {
    const r = await fetch('/api/personalities');
    const data = await r.json();
    _personalities = data.personalities || [];
    const sel = document.getElementById('personalitySelect');
    while (sel.options.length > 1) sel.remove(1);
    for (const p of _personalities) {
      const opt = document.createElement('option');
      opt.value = p.name;
      opt.textContent = p.name;
      sel.appendChild(opt);
    }
    if (data.active) sel.value = data.active;
  } catch(e) { console.warn('Personalities:', e); }
}

async function activatePersonality() {
  const sel = document.getElementById('personalitySelect');
  const name = sel.value || null;
  try {
    const r = await fetch('/api/personalities/set', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({name})
    });
    const d = await r.json();
    if (d.ok) toast(name ? `◈ PERSONA: ${name}` : '◈ NOVA DEFAULT RESTORED');
    else toast('❌ ' + (d.error || 'Personality error'), true);
  } catch(e) { toast('❌ Personality error', true); }
}

loadPersonalities();

async function loadHistory() {
  try {
    const r = await fetch('/api/history');
    const hist = await r.json();
    const newLast = hist.length ? hist[hist.length-1].content : '';
    if ((hist.length !== lastCount || newLast !== lastContent) && !isThinking) {
      lastContent = newLast;
      const prev = lastCount;
      const sp = conv.scrollHeight - conv.scrollTop;
      conv.innerHTML = '';
      hist.forEach(m => addMessage(m.role, m.content, false));
      lastCount = hist.length;
      if (!userScrolled) conv.scrollTop = conv.scrollHeight;
      else conv.scrollTop = conv.scrollHeight - sp;
      if (hist.length > prev) {
        const last = hist[hist.length-1];
        if (last.role === 'assistant') playTTS(last.content);
      }
    }
  } catch(e) {}
}

const dropZone = document.getElementById('drop-zone');
const fileInput = document.getElementById('file-input');
const filePrev = document.getElementById('file-prev');

['dragenter','dragover','dragleave','drop'].forEach(ev => {
  dropZone.addEventListener(ev, e => { e.preventDefault(); e.stopPropagation(); });
  document.body.addEventListener(ev, e => { e.preventDefault(); e.stopPropagation(); });
});
['dragenter','dragover'].forEach(ev => dropZone.addEventListener(ev, () => dropZone.classList.add('drag-over')));
['dragleave','drop'].forEach(ev => dropZone.addEventListener(ev, () => dropZone.classList.remove('drag-over')));
dropZone.addEventListener('drop', e => uploadFiles(e.dataTransfer.files));
dropZone.addEventListener('click', () => fileInput.click());
fileInput.addEventListener('change', e => { uploadFiles(e.target.files); fileInput.value = ''; });

async function uploadFiles(files) {
  if (!files.length) return;
  const fd = new FormData();
  for (const f of files) fd.append('files', f);
  toast('⏳ UPLOADING...');
  try {
    const r = await fetch('/api/upload', {method:'POST', body:fd});
    const d = await r.json();
    if (d.status === 'success') {
      d.files.forEach(f => pendingFiles.push(f));
      showFilePreviews(d.files);
      toast(`✓ ${d.files.length} file(s) staged`);
    }
  } catch { toast('❌ Upload failed', true); }
}

function showFilePreviews(files) {
  files.forEach((f, i) => {
    const idx = pendingFiles.length - files.length + i;
    const icon = f.original_name.match(/\.pdf$/i) ? '📄' :
                 f.original_name.match(/\.(png|jpg|jpeg|gif|webp)$/i) ? '🖼' : '📎';
    const el = document.createElement('div');
    el.className = 'fpi';
    el.innerHTML = `<span>${icon}</span><span>${f.original_name} (${(f.size/1024).toFixed(1)}KB)</span><button onclick="removePending(${idx},this.parentElement)">✕</button>`;
    filePrev.appendChild(el);
  });
}

function removePending(idx, el) {
  pendingFiles[idx] = null;
  el.remove();
}

function buildFilePayload() {
  const active = pendingFiles.filter(f => f !== null);
  if (!active.length) return null;
  let p = `I've attached ${active.length} file(s):\n\n`;
  active.forEach(f => {
    p += `**File: ${f.original_name}**\nSize: ${(f.size/1024).toFixed(1)} KB\nContent:\n\`\`\`\n${f.preview}\n\`\`\`\n\n`;
  });
  pendingFiles = [];
  filePrev.innerHTML = '';
  return p;
}

msgIn.addEventListener('input', function() {
  this.style.height = 'auto';
  this.style.height = Math.min(this.scrollHeight, 130) + 'px';
  const n = this.value.length;
  const ci = document.getElementById('char-info');
  ci.textContent = n > 200 ? n + ' / 50000' : '';
  ci.style.color = n > 45000 ? 'var(--coral)' : 'var(--dimteal)';
});
msgIn.addEventListener('keydown', e => {
  if (e.key === 'Enter' && !e.shiftKey) { e.preventDefault(); sendMessage(); }
});
conv.addEventListener('scroll', () => {
  const atBot = conv.scrollHeight - conv.scrollTop - conv.clientHeight < 60;
  userScrolled = !atBot;
  if (atBot) userScrolled = false;
});

// ============================================================
//  INITIALIZATION
// ============================================================
const starCanvas = document.getElementById('star-canvas');
const starfield = new StarfieldCanvas(starCanvas);

const dodecCanvas = document.getElementById('dodec-canvas');
const dodec = new DodecahedronAnim(dodecCanvas);

// Initialize ambient hum on first user interaction
document.addEventListener('click', () => {
  if (ambientEnabled && !ambientContext) initAmbient();
}, { once: true });

loadHistory();
setInterval(loadHistory, 1500);
setTimeout(() => msgIn.focus(), 400);
toast('🖖 NOVA HYPERION CORE ONLINE');

// Sync ambient state with server
fetch('/api/state').then(r => r.json()).then(state => {
  if (state.ambient_hum === false) toggleAmbient();
});
</script>
</body>
</html>
"""

# The NovaWebServer class and main application integration remain the same
# as in the original web interface code